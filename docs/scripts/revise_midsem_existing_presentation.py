from pptx import Presentation

from build_midsem_presentation_v2 import (
    BLUE,
    CREAM,
    GREEN,
    NAVY,
    OUT_NOTES,
    OUT_PPTX,
    PALE_GOLD,
    ROSE,
    SLATE,
    SLATE_BG,
    WHITE,
    BODY_FONT,
    TITLE_FONT,
    bullets,
    card,
    number,
    rect,
    set_bg,
    style,
    text,
    title,
)
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches


def delete_slide(prs, index):
    slide_id_list = prs.slides._sldIdLst
    slides = list(slide_id_list)
    slide_id = slides[index]
    prs.part.drop_rel(slide_id.rId)
    slide_id_list.remove(slide_id)


def clear_slide(slide):
    for shape in list(slide.shapes):
        shape.element.getparent().remove(shape.element)


def reset_number(slide, value, dark=False):
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip().isdigit():
            shape.text = str(value)
            paragraph = shape.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.RIGHT
            style(paragraph, 10, PALE_GOLD if dark else SLATE, False, BODY_FONT)
            return
    number(slide, value, dark=dark)


def rebuild_solution_slide(slide):
    clear_slide(slide)
    set_bg(slide, CREAM)
    title(slide, "Solution Overview", "A web-based system with separate candidate and employer interfaces", "Solution")
    text(
        slide,
        Inches(0.78),
        Inches(1.55),
        Inches(11.7),
        Inches(0.42),
        "Job Hunt is an initial matching and recommendation stage for computing-related roles, not a full hiring platform.",
        15,
        NAVY,
        True,
    )
    card(
        slide,
        Inches(0.72),
        Inches(2.1),
        Inches(3.75),
        Inches(3.5),
        "Candidate Web Interface",
        [
            "Upload a resume.",
            "View ranked job matches.",
            "View match details.",
            "Possible explanation features remain optional.",
        ],
        BLUE,
    )
    card(
        slide,
        Inches(4.8),
        Inches(2.1),
        Inches(3.75),
        Inches(3.5),
        "Core Matching Engine",
        [
            "Extract competencies from resumes and job descriptions.",
            "Map extracted signals to O*NET descriptors.",
            "Determine whether required competencies are satisfied.",
            "Rank useful results for both sides.",
        ],
        GREEN,
    )
    card(
        slide,
        Inches(8.88),
        Inches(2.1),
        Inches(3.75),
        Inches(3.5),
        "Employer Web Interface",
        [
            "Create a job posting.",
            "View ranked candidates.",
            "View match details.",
            "Preferred vs required controls remain stretch features.",
        ],
        PALE_GOLD,
    )


def rebuild_progress_slide(slide):
    clear_slide(slide)
    set_bg(slide, WHITE)
    title(slide, "Progress on Target and Well Organised", "Current progress, remaining work, and work division", "Progress")
    card(
        slide,
        Inches(0.62),
        Inches(1.68),
        Inches(3.95),
        Inches(4.95),
        "What has been done",
        [
            "Problem definition and scope have been clarified.",
            "The O*NET-based methodology has been outlined.",
            "Functional requirements have been drafted.",
            "Sitemap and use case diagrams have been prepared.",
            "Work areas across the team are already assigned.",
        ],
        BLUE,
    )
    card(
        slide,
        Inches(4.7),
        Inches(1.68),
        Inches(3.95),
        Inches(4.95),
        "What remains to be done",
        [
            "Implement the system as a web application.",
            "Complete frontend and backend integration.",
            "Finish database design and setup.",
            "Prepare datasets and carry out evaluation.",
            "Use Visual Paradigm and StarUML for modelling and documentation support.",
        ],
        PALE_GOLD,
    )
    card(
        slide,
        Inches(8.78),
        Inches(1.68),
        Inches(3.95),
        Inches(4.95),
        "Work division",
        [
            "Alwyn Sterling - backend development and database implementation",
            "Dominic Scott - backend development and database implementation",
            "Ron-Hugh Walters - documentation, modelling, data cleaning, and parsing",
            "Leah-Jay Holness - user interface design",
            "Tianna Bassaragh - user interface design",
        ],
        GREEN,
    )


def write_notes():
    notes = """# Job Hunt Midsemester Presentation Notes

## Slide 1 - Title slide
- Project name, team name, course name, and team members only.

## Slide 2 - Problem background and context
- Recruitment is inefficient and often opaque.
- The project focuses on computing-related roles.
- The system is an initial matching stage, not a full hiring platform.

## Slide 3 - Solution overview
- Present the project clearly as a web application.
- Separate the candidate interface, the core matching engine, and the employer interface.
- Keep stretch features framed as optional.

## Slide 4 - Objectives clearly defined
- Core deliverables, what is not in scope, and what may be added if time permits.

## Slide 5 - Functional requirements overview
- Focus on the key required behaviours only.

## Slide 6 - System flow and methodology
- Collect, extract, map, qualify, and rank.
- Mention evaluation metrics briefly.

## Slide 7 - Sitemap
- Explain public, candidate, and employer areas.

## Slide 8 - Use case diagram
- Highlight the main actors and interactions.

## Slide 9 - Technical challenges and why the project is non-trivial
- Extraction, O*NET mapping, defensible qualification, and evaluation.

## Slide 10 - Progress on target and well organised
- Keep the done, remaining, and work division sections concise.

## Slide 11 - Knowledge from the curriculum and next steps
- Close on curriculum relevance and the realistic next implementation steps.
"""
    OUT_NOTES.write_text(notes, encoding="utf-8")


def main():
    prs = Presentation(str(OUT_PPTX))
    delete_slide(prs, 1)
    rebuild_solution_slide(prs.slides[2])
    rebuild_progress_slide(prs.slides[9])
    for i, slide in enumerate(prs.slides, 1):
        reset_number(slide, i, dark=(i == 1))
    prs.save(str(OUT_PPTX))
    write_notes()
    print(f"Revised {OUT_PPTX}")
    print(f"Updated {OUT_NOTES}")


if __name__ == "__main__":
    main()
