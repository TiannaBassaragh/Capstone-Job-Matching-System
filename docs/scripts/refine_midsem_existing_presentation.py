from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from build_midsem_presentation_v2 import (
    BLUE,
    CREAM,
    DOCS,
    GREEN,
    NAVY,
    OUT_NOTES,
    OUT_PPTX,
    PALE_GOLD,
    ROSE,
    SLATE,
    SLATE_BG,
    WHITE,
    BODY_FONT,
    card,
    image_box,
    number,
    rect,
    set_bg,
    style,
    text,
    title,
)
from pptx.enum.text import PP_ALIGN


ERD = DOCS / "Schema" / "ERD for schema.png"


def clear_slide(slide):
    for shape in list(slide.shapes):
        shape.element.getparent().remove(shape.element)


def slide_title(slide):
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            return shape.text.strip().splitlines()[0]
    return ""


def find_slide_index(prs, wanted):
    for i, slide in enumerate(prs.slides):
        if slide_title(slide) == wanted:
            return i
    raise ValueError(f"Slide not found: {wanted}")


def move_slide(prs, old_index, new_index):
    slide_id_list = prs.slides._sldIdLst
    slides = list(slide_id_list)
    slide = slides[old_index]
    slide_id_list.remove(slide)
    slide_id_list.insert(new_index, slide)


def reset_number(slide, value, dark=False):
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip().isdigit():
            shape.text = str(value)
            paragraph = shape.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.RIGHT
            style(paragraph, 10, PALE_GOLD if dark else SLATE, False, BODY_FONT)
            return
    number(slide, value, dark=dark)


def rebuild_solution_slide(slide):
    clear_slide(slide)
    set_bg(slide, CREAM)
    title(slide, "Solution Overview", "A web-based Job Matching and Qualification System for computing-related roles", "Solution")
    text(
        slide,
        Inches(0.78),
        Inches(1.55),
        Inches(11.7),
        Inches(0.42),
        "The system operates as an initial matching and recommendation stage, not a full hiring platform.",
        15,
        NAVY,
        True,
    )
    card(slide, Inches(0.72), Inches(2.1), Inches(3.75), Inches(3.45), "Candidate Web Interface", ["Upload resume", "View ranked job matches", "View match details"], BLUE)
    card(slide, Inches(4.8), Inches(2.1), Inches(3.75), Inches(3.45), "Core Matching Engine", ["Extract competencies", "Map competencies to O*NET descriptors", "Determine qualification", "Rank results"], GREEN)
    card(slide, Inches(8.88), Inches(2.1), Inches(3.75), Inches(3.45), "Employer Web Interface", ["Create job posting", "View ranked candidates", "View match details"], PALE_GOLD)
    text(
        slide,
        Inches(0.85),
        Inches(6.05),
        Inches(11.4),
        Inches(0.28),
        "Possible explanation, gap-report, weighting, and advanced filtering features remain stretch additions rather than guaranteed deliverables.",
        11.5,
        SLATE,
    )


def rebuild_progress_slide(slide):
    clear_slide(slide)
    set_bg(slide, WHITE)
    title(slide, "Progress on Target and Well Organised", "Progress so far has focused on planning, modelling, and implementation preparation", "Progress")
    card(slide, Inches(0.62), Inches(1.7), Inches(3.95), Inches(4.9), "What has been done so far", ["Problem definition and scope have been clarified.", "The O*NET-based methodology has been outlined.", "Functional requirements have been drafted.", "Use case and sitemap artefacts have been prepared.", "Initial presentation structure has been organised around the rubric."], BLUE)
    card(slide, Inches(4.7), Inches(1.7), Inches(3.95), Inches(4.9), "What remains to be done", ["Implement the system as a web application.", "Complete frontend and backend integration.", "Finish database design and setup.", "Prepare datasets and carry out evaluation.", "Use Visual Paradigm and StarUML for modelling and documentation support."], PALE_GOLD)
    card(slide, Inches(8.78), Inches(1.7), Inches(3.95), Inches(4.9), "Work division", ["Alwyn Sterling - backend development and database implementation", "Dominic Scott - backend development and database implementation", "Ron-Hugh Walters - documentation, modelling, data cleaning, and parsing", "Leah-Jay Holness - user interface design", "Tianna Bassaragh - user interface design"], GREEN)


def rebuild_use_case_slide(slide):
    clear_slide(slide)
    set_bg(slide, WHITE)
    title(slide, "Use Case Diagram", "Progress artefact for the planned user interactions", "Progress")
    rect(slide, Inches(0.62), Inches(1.55), Inches(8.65), Inches(5.25), SLATE_BG, radius=True)
    image_box(slide, DOCS / "use_case_diagram" / "UseCaseDiagram7.png", Inches(0.85), Inches(1.78), Inches(8.2), Inches(4.75))
    card(slide, Inches(9.55), Inches(1.78), Inches(2.95), Inches(3.25), "Why it matters", ["This artefact clarifies the main interactions between Candidates, Employers, and the system.", "It helps refine what each actor can submit, view, and inspect."], BLUE)
    text(slide, Inches(0.88), Inches(6.56), Inches(8.15), Inches(0.25), "Caption: This artefact clarifies the main interactions between Candidates, Employers, and the system.", 11.5, SLATE)


def rebuild_sitemap_slide(slide):
    clear_slide(slide)
    set_bg(slide, CREAM)
    title(slide, "Sitemap", "Progress artefact for the planned web application structure", "Progress")
    rect(slide, Inches(0.62), Inches(1.55), Inches(8.65), Inches(5.25), SLATE_BG, radius=True)
    image_box(slide, DOCS / "Sitemap" / "Real Site Map.png", Inches(0.85), Inches(1.78), Inches(8.2), Inches(4.75))
    card(slide, Inches(9.55), Inches(1.78), Inches(2.95), Inches(3.25), "Why it matters", ["This artefact plans the front end structure of the web application for candidate and employer workflows.", "It separates public pages from candidate and employer dashboard paths."], PALE_GOLD)
    text(slide, Inches(0.88), Inches(6.56), Inches(8.15), Inches(0.25), "Caption: This artefact plans the front end structure of the web application for candidate and employer workflows.", 11.5, SLATE)


def add_erd_slide(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    title(slide, "Database Design Overview", "Progress artefact for the proposed persistence layer", "Progress")
    rect(slide, Inches(0.62), Inches(1.55), Inches(8.65), Inches(5.25), SLATE_BG, radius=True)
    image_box(slide, ERD, Inches(0.85), Inches(1.78), Inches(8.2), Inches(4.75))
    card(
        slide,
        Inches(9.55),
        Inches(1.72),
        Inches(2.95),
        Inches(3.55),
        "Schema highlights",
        [
            "Supports candidate and employer data in the web application.",
            "Stores job posts, resumes, and competencies as core entities.",
            "Uses linking tables for many-to-many competency relationships.",
            "Stores match outcomes, qualification status, and match scores.",
        ],
        GREEN,
    )
    text(
        slide,
        Inches(0.88),
        Inches(6.56),
        Inches(8.15),
        Inches(0.25),
        "Caption: This artefact defines the proposed database structure needed to support resumes, job postings, competencies, and match results.",
        11.5,
        SLATE,
    )


def write_notes():
    notes = """# Job Hunt Midsemester Presentation Notes

## Slide 1 - Title slide
- Project name, team name, course name, and team members only.

## Slide 2 - Problem background and context
- Introduce the recruitment problem and project scope.

## Slide 3 - Solution overview
- Present the project clearly as a web application with candidate and employer interfaces.
- The core engine performs extraction, O*NET mapping, qualification, and ranking.

## Slide 4 - Objectives clearly defined
- Keep must-do features, out-of-scope items, and stretch goals distinct.

## Slide 5 - Functional requirements overview
- Focus on the most important required behaviours only.

## Slide 6 - System flow and methodology
- Explain the flow from input through ranking and mention evaluation metrics briefly.

## Slide 7 - Technical challenges and why the project is non-trivial
- Extraction, O*NET mapping, defensible qualification, and evaluation.

## Slide 8 - Progress on target and well organised
- Emphasise planning, modelling, and implementation preparation.

## Slide 9 - Use case diagram
- Frame the diagram as a progress artefact for interaction design.

## Slide 10 - Sitemap
- Frame the diagram as a progress artefact for front-end planning.

## Slide 11 - Database design overview
- Frame the ERD as a progress artefact for the planned schema.

## Slide 12 - Knowledge from the curriculum and next steps
- Close on coursework relevance and the realistic next implementation steps.
"""
    OUT_NOTES.write_text(notes, encoding="utf-8")


def main():
    prs = Presentation(str(OUT_PPTX))

    rebuild_solution_slide(prs.slides[2])
    rebuild_progress_slide(prs.slides[9])
    rebuild_use_case_slide(prs.slides[7])
    rebuild_sitemap_slide(prs.slides[6])
    add_erd_slide(prs)

    move_slide(prs, find_slide_index(prs, "Technical Challenges and Why the Project Is Non-Trivial"), 6)
    move_slide(prs, find_slide_index(prs, "Progress on Target and Well Organised"), 7)
    move_slide(prs, find_slide_index(prs, "Use Case Diagram"), 8)
    move_slide(prs, find_slide_index(prs, "Sitemap"), 9)
    move_slide(prs, find_slide_index(prs, "Database Design Overview"), 10)

    for i, slide in enumerate(prs.slides, 1):
        reset_number(slide, i, dark=(i == 1))

    prs.save(str(OUT_PPTX))
    write_notes()
    print(f"Refined {OUT_PPTX}")
    print(f"Updated {OUT_NOTES}")


if __name__ == "__main__":
    main()
