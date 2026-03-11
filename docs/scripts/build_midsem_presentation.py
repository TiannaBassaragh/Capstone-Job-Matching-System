from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as ShapeType
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "docs"
OUT_PPTX = OUT_DIR / "job-hunt-midsemester-presentation.pptx"
OUT_NOTES = OUT_DIR / "job-hunt-midsemester-notes.md"
SITE_MAP = ROOT / "docs" / "Real Site Map.png"
USE_CASE = ROOT / "docs" / "use_case_diagram" / "UseCaseDiagram7.png"

NAVY = RGBColor(17, 33, 62)
SLATE = RGBColor(45, 62, 84)
GOLD = RGBColor(224, 180, 74)
TEAL = RGBColor(54, 138, 127)
CREAM = RGBColor(247, 243, 235)
WHITE = RGBColor(255, 255, 255)
INK = RGBColor(31, 42, 55)
BLUE = RGBColor(225, 235, 247)
GREEN = RGBColor(228, 242, 236)
PALE_GOLD = RGBColor(247, 235, 199)
ROSE = RGBColor(246, 225, 220)
SLATE_BG = RGBColor(234, 238, 243)

TITLE_FONT = "Georgia"
BODY_FONT = "Aptos"


def ensure_dirs():
    OUT_DIR.mkdir(exist_ok=True)


def style(paragraph, size, color, bold=False, font=BODY_FONT):
    if not paragraph.runs:
        paragraph.add_run()
    for run in paragraph.runs:
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, left, top, width, height, fill, radius=False):
    kind = ShapeType.ROUNDED_RECTANGLE if radius else ShapeType.RECTANGLE
    shape = slide.shapes.add_shape(kind, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    return shape


def text(slide, left, top, width, height, value, size=18, color=INK, bold=False, font=BODY_FONT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Pt(6)
    frame.margin_right = Pt(6)
    frame.margin_top = Pt(4)
    frame.margin_bottom = Pt(4)
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.text = value
    p.alignment = align
    style(p, size, color, bold=bold, font=font)
    return box


def bullets(slide, left, top, width, height, items, size=15, color=INK):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Pt(8)
    frame.margin_right = Pt(8)
    frame.margin_top = Pt(4)
    frame.margin_bottom = Pt(4)
    frame.clear()
    for i, item in enumerate(items):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = item
        p.bullet = True
        p.space_after = Pt(7)
        style(p, size, color)
    return box


def title(slide, value, subtitle, section, dark=False):
    rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.34), GOLD)
    title_color = WHITE if dark else NAVY
    sub_color = PALE_GOLD if dark else TEAL
    text(slide, Inches(0.7), Inches(0.55), Inches(9.5), Inches(0.75), value, 25, title_color, True, TITLE_FONT)
    text(slide, Inches(0.72), Inches(1.2), Inches(9.7), Inches(0.38), subtitle, 11, sub_color, True)
    tag = rect(slide, Inches(10.72), Inches(0.62), Inches(1.8), Inches(0.36), TEAL, radius=True)
    p = tag.text_frame.paragraphs[0]
    p.text = section
    p.alignment = PP_ALIGN.CENTER
    style(p, 11, WHITE, True)


def card(slide, left, top, width, height, heading, items, fill):
    rect(slide, left, top, width, height, fill, radius=True)
    text(slide, left + Inches(0.12), top + Inches(0.1), width - Inches(0.24), Inches(0.34), heading, 14, NAVY, True, TITLE_FONT)
    bullets(slide, left + Inches(0.08), top + Inches(0.46), width - Inches(0.16), height - Inches(0.54), items, 11.5)


def step(slide, left, top, heading, body, fill):
    rect(slide, left, top, Inches(1.45), Inches(1.1), fill, radius=True)
    text(slide, left + Inches(0.1), top + Inches(0.08), Inches(1.25), Inches(0.28), heading, 13, NAVY, True, TITLE_FONT, PP_ALIGN.CENTER)
    text(slide, left + Inches(0.08), top + Inches(0.38), Inches(1.28), Inches(0.55), body, 10.5, INK, False, BODY_FONT, PP_ALIGN.CENTER)


def chevron(slide, left, top):
    shape = slide.shapes.add_shape(ShapeType.CHEVRON, left, top, Inches(0.42), Inches(0.54))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.color.rgb = GOLD


def image_box(slide, path, left, top, width, height):
    with Image.open(path) as image:
        iw, ih = image.size
    box_ratio = float(width) / float(height)
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        draw_width = width
        draw_height = int(width / img_ratio)
    else:
        draw_height = height
        draw_width = int(height * img_ratio)
    draw_left = left + int((width - draw_width) / 2)
    draw_top = top + int((height - draw_height) / 2)
    slide.shapes.add_picture(str(path), draw_left, draw_top, width=draw_width, height=draw_height)


def number(slide, value, dark=False):
    color = PALE_GOLD if dark else SLATE
    text(slide, Inches(12.52), Inches(7.0), Inches(0.4), Inches(0.2), str(value), 10, color, False, BODY_FONT, PP_ALIGN.RIGHT)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slide_no = 1

    slide = prs.slides.add_slide(blank)
    set_bg(slide, NAVY)
    rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.34), GOLD)
    rect(slide, Inches(0), Inches(6.96), Inches(13.333), Inches(0.54), TEAL)
    rect(slide, Inches(8.55), Inches(1.2), Inches(3.82), Inches(4.45), PALE_GOLD, radius=True)
    text(slide, Inches(0.78), Inches(1.0), Inches(5.5), Inches(0.35), "COMP3901 Midsemester Presentation", 14, PALE_GOLD, True)
    text(slide, Inches(0.78), Inches(1.55), Inches(5.5), Inches(0.9), "Job Hunt", 30, WHITE, True, TITLE_FONT)
    text(
        slide,
        Inches(0.8),
        Inches(2.42),
        Inches(6.25),
        Inches(1.0),
        "A job matching and qualification system for computing-related roles",
        20,
        WHITE,
    )
    text(slide, Inches(0.8), Inches(3.55), Inches(2.0), Inches(0.3), "Team Name", 12, PALE_GOLD, True)
    text(slide, Inches(0.8), Inches(3.86), Inches(3.0), Inches(0.5), "Free Labour", 22, WHITE, True, TITLE_FONT)
    text(slide, Inches(8.92), Inches(1.52), Inches(2.8), Inches(0.3), "Presented by", 12, NAVY, True)
    bullets(
        slide,
        Inches(8.86),
        Inches(1.9),
        Inches(3.0),
        Inches(2.75),
        [
            "Alwyn Sterling",
            "Dominic Scott",
            "Leah-Jay Holness",
            "Ron-Hugh Walters",
            "Tianna Bassaragh",
        ],
        14,
    )
    text(slide, Inches(8.92), Inches(5.15), Inches(2.8), Inches(0.3), "Midsemester focus", 11, TEAL, True)
    text(
        slide,
        Inches(8.92),
        Inches(5.42),
        Inches(2.9),
        Inches(0.8),
        "Problem definition, solution approach, scope, and progress to date",
        14,
        NAVY,
    )
    number(slide, slide_no, dark=True)
    slide_no += 1

    slide = prs.slides.add_slide(blank)
    set_bg(slide, CREAM)
    title(slide, "Objectives Clearly Defined", "Background and context for the project", "Context")
    bullets(
        slide,
        Inches(0.72),
        Inches(1.55),
        Inches(6.45),
        Inches(4.9),
        [
            "Modern recruitment is often inefficient and opaque for both employers and applicants.",
            "Employers sift through many applications of varying relevance, while applicants often receive little feedback.",
            "This causes wasted effort, missed qualified candidates, and weak transparency in the matching process.",
            "The project formalizes comparison between job requirements and candidate competencies into a structured computational process.",
            "Scope is limited to computing-related roles so the problem stays manageable and evaluation stays realistic.",
            "The system is an initial matching stage, not a full hiring workflow and not a replacement for recruiters.",
        ],
        16,
    )
    card(
        slide,
        Inches(7.55),
        Inches(1.7),
        Inches(5.05),
        Inches(1.12),
        "Why this matters",
        [
            "Qualified candidates can be missed.",
            "Employers spend time filtering weak matches.",
            "Applicants get little explanation.",
        ],
        PALE_GOLD,
    )
    card(
        slide,
        Inches(7.55),
        Inches(3.02),
        Inches(5.05),
        Inches(1.12),
        "Project boundary",
        [
            "We focus on resumes, job descriptions, and competency-based qualification.",
            "We do not model final hiring outcomes.",
        ],
        BLUE,
    )
    card(
        slide,
        Inches(7.55),
        Inches(4.34),
        Inches(5.05),
        Inches(1.22),
        "Core framing",
        [
            "This is a structured decision-support problem.",
            "The goal is a clearer first-pass judgement of suitability.",
        ],
        GREEN,
    )
    number(slide, slide_no)
    slide_no += 1

    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    title(slide, "Proposed Solution", "How Job Hunt is meant to work", "Approach")
    bullets(
        slide,
        Inches(0.72),
        Inches(1.55),
        Inches(5.55),
        Inches(5.0),
        [
            "Accept resumes from candidates and job descriptions from employers.",
            "Extract relevant competencies from both forms of unstructured text.",
            "Map extracted competencies into a shared O*NET-based competency space.",
            "Determine whether a candidate is qualified for a role, then rank by suitability.",
            "Reduce false recommendations and missed valid matches.",
            "Approximate the first-pass judgement of a skilled recruiter in a more structured way.",
        ],
        16,
    )
    step(slide, Inches(6.55), Inches(1.95), "Input", "Resumes and job descriptions", PALE_GOLD)
    chevron(slide, Inches(8.03), Inches(2.22))
    step(slide, Inches(8.52), Inches(1.95), "Extract", "Pull out competencies from text", GREEN)
    chevron(slide, Inches(10.0), Inches(2.22))
    step(slide, Inches(10.48), Inches(1.95), "Map", "Align to shared O*NET descriptors", BLUE)
    text(slide, Inches(9.78), Inches(3.45), Inches(2.25), Inches(0.3), "Then", 13, TEAL, True, BODY_FONT, PP_ALIGN.CENTER)
    step(slide, Inches(9.7), Inches(3.78), "Decide", "Qualify and rank matches", ROSE)
    card(
        slide,
        Inches(6.55),
        Inches(5.05),
        Inches(5.75),
        Inches(1.0),
        "Central objective",
        [
            "Find candidates who genuinely satisfy requirements.",
            "Avoid surfacing candidates who do not meet required competencies.",
        ],
        SLATE_BG,
    )
    number(slide, slide_no)
    slide_no += 1

    slide = prs.slides.add_slide(blank)
    set_bg(slide, CREAM)
    title(slide, "Project Scope", "What we will do, what we will not do, and what may be added later", "Scope")
    card(
        slide,
        Inches(0.62),
        Inches(1.55),
        Inches(4.05),
        Inches(4.95),
        "We will definitely do",
        [
            "Accept resumes and job descriptions as input.",
            "Extract competencies from unstructured text.",
            "Map extracted competencies to O*NET descriptors.",
            "Generate structured competency profiles for jobs and candidates.",
            "Determine whether required competencies are satisfied.",
            "Rank qualified candidates and relevant jobs.",
        ],
        GREEN,
    )
    card(
        slide,
        Inches(4.89),
        Inches(1.55),
        Inches(4.0),
        Inches(4.95),
        "We will definitely not do",
        [
            "Build a complete hiring platform.",
            "Handle interviews, onboarding, or final hiring decisions.",
            "Predict actual job performance.",
            "Guarantee real-world hiring accuracy from resume text alone.",
        ],
        ROSE,
    )
    card(
        slide,
        Inches(9.08),
        Inches(1.55),
        Inches(3.65),
        Inches(4.95),
        "We might do if time permits",
        [
            "Mark competencies as required or preferred.",
            "Adjust ranking weights within limits.",
            "Generate qualification explanations.",
            "Produce competency gap reports.",
            "Add user result filters.",
        ],
        PALE_GOLD,
    )
    number(slide, slide_no)
    slide_no += 1

    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    title(
        slide,
        "Technical Issues and Why the Project Is Non-Trivial",
        "The hard parts are modelling and evaluation, not just the interface",
        "Challenge",
    )
    card(
        slide,
        Inches(0.72),
        Inches(1.7),
        Inches(2.85),
        Inches(2.0),
        "Messy text",
        [
            "Resumes and job descriptions are vague, inconsistent, and full of implied skills.",
            "Competency extraction must work on noisy language, not clean forms.",
        ],
        BLUE,
    )
    card(
        slide,
        Inches(3.78),
        Inches(1.7),
        Inches(2.85),
        Inches(2.0),
        "Defining qualified",
        [
            "Qualification is multi-dimensional, not a single score.",
            "Multiple signals must be combined without becoming arbitrary.",
        ],
        GREEN,
    )
    card(
        slide,
        Inches(6.84),
        Inches(1.7),
        Inches(2.85),
        Inches(2.0),
        "High-dimensional space",
        [
            "The system assumes useful structure in a competency space.",
            "Linearity may not always hold, and scale still matters.",
        ],
        PALE_GOLD,
    )
    card(
        slide,
        Inches(9.9),
        Inches(1.7),
        Inches(2.75),
        Inches(2.0),
        "Evaluation",
        [
            "Ground truth for who is truly qualified is hard to obtain.",
            "Testing must be credible, staged, and metric-driven.",
        ],
        ROSE,
    )
    rect(slide, Inches(0.72), Inches(4.28), Inches(11.9), Inches(1.5), NAVY, radius=True)
    text(slide, Inches(0.98), Inches(4.5), Inches(3.0), Inches(0.3), "Evaluation plan from the draft", 12, PALE_GOLD, True)
    bullets(
        slide,
        Inches(0.95),
        Inches(4.82),
        Inches(11.3),
        Inches(0.8),
        [
            "Stage 1: qualification decision using precision, recall, and F1.",
            "Stage 2: ranking quality using Precision@K, nDCG@K, and MRR if time permits.",
        ],
        15,
        WHITE,
    )
    number(slide, slide_no)
    slide_no += 1

    slide = prs.slides.add_slide(blank)
    set_bg(slide, CREAM)
    title(slide, "Knowledge From the Curriculum Required", "Why this is a capstone-level problem", "Curriculum")
    course_cards = [
        ("COMP2171 + COMP2140", ["Documentation, requirements analysis, system modelling, and design artifacts."], BLUE),
        ("COMP3161", ["Database design for resumes, job posts, competency profiles, and match results."], GREEN),
        ("COMP3162", ["Data cleaning, preprocessing, modelling support, and evaluation work."], PALE_GOLD),
        ("COMP1127", ["Python foundations for backend logic, parsing, and experimentation."], ROSE),
        ("INFO2180", ["Frontend design and implementation for employer and candidate interaction."], SLATE_BG),
        ("COMP3220", ["Intelligent decision-making for extraction, matching, and qualification."], BLUE),
        ("COMP2211", ["Optimization thinking for ranking and balancing multiple decision factors."], GREEN),
    ]
    positions = [
        (Inches(0.72), Inches(1.65)),
        (Inches(3.95), Inches(1.65)),
        (Inches(7.18), Inches(1.65)),
        (Inches(10.41), Inches(1.65)),
        (Inches(2.05), Inches(4.1)),
        (Inches(5.28), Inches(4.1)),
        (Inches(8.51), Inches(4.1)),
    ]
    for (heading, items, fill), (left, top) in zip(course_cards, positions):
        card(slide, left, top, Inches(2.45), Inches(1.78), heading, items, fill)
    text(
        slide,
        Inches(0.88),
        Inches(6.45),
        Inches(11.8),
        Inches(0.35),
        "Taken together, the project draws on software engineering, databases, web development, AI, data science, programming, and optimization.",
        13,
        NAVY,
        True,
    )
    number(slide, slide_no)
    slide_no += 1

    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    title(slide, "Progress on Target", "Who is doing what right now", "Progress")
    card(slide, Inches(0.72), Inches(1.7), Inches(3.0), Inches(1.2), "Alwyn Sterling", ["Backend development", "Database implementation"], BLUE)
    card(slide, Inches(0.72), Inches(3.05), Inches(3.0), Inches(1.2), "Dominic Scott", ["Backend development", "Database implementation"], GREEN)
    card(
        slide,
        Inches(0.72),
        Inches(4.4),
        Inches(3.0),
        Inches(1.45),
        "Ron-Hugh Walters",
        ["Documentation", "System modelling", "Data cleaning", "Data parsing"],
        PALE_GOLD,
    )
    card(slide, Inches(3.95), Inches(1.7), Inches(3.0), Inches(1.2), "Leah-Jay Holness", ["User interface design"], ROSE)
    card(slide, Inches(3.95), Inches(3.05), Inches(3.0), Inches(1.2), "Tianna Bassaragh", ["User interface design"], SLATE_BG)
    rect(slide, Inches(7.42), Inches(1.75), Inches(5.0), Inches(4.05), NAVY, radius=True)
    text(slide, Inches(7.72), Inches(2.0), Inches(4.2), Inches(0.3), "What this shows for midsemester", 13, PALE_GOLD, True)
    bullets(
        slide,
        Inches(7.65),
        Inches(2.34),
        Inches(4.3),
        Inches(3.05),
        [
            "The work is split across backend, database, UI, documentation, modelling, and data preparation.",
            "The team already has a defined structure rather than an undefined group effort.",
            "That makes the remaining implementation plan easier to organize and defend.",
        ],
        15,
        WHITE,
    )
    number(slide, slide_no)
    slide_no += 1

    slide = prs.slides.add_slide(blank)
    set_bg(slide, CREAM)
    title(slide, "What Remains to Be Done", "Main implementation and evaluation work still ahead", "Next")
    card(
        slide,
        Inches(0.8),
        Inches(1.6),
        Inches(7.0),
        Inches(1.02),
        "1. Frontend",
        ["Candidate and employer interfaces, page flow, and user interaction."],
        BLUE,
    )
    card(
        slide,
        Inches(0.8),
        Inches(2.9),
        Inches(7.0),
        Inches(1.02),
        "2. Backend + Database",
        ["API endpoints, document processing flow, and storage for jobs, resumes, profiles, and matches."],
        GREEN,
    )
    card(
        slide,
        Inches(0.8),
        Inches(4.2),
        Inches(7.0),
        Inches(1.02),
        "3. Implement the matching pipeline",
        ["Resume parsing, job parsing, mapping, qualification, and ranking logic."],
        PALE_GOLD,
    )
    card(
        slide,
        Inches(0.8),
        Inches(5.5),
        Inches(7.0),
        Inches(1.02),
        "4. Validate the system",
        ["Prepare datasets, clean data, integrate evaluation, and test the pipeline."],
        ROSE,
    )
    card(
        slide,
        Inches(8.35),
        Inches(1.82),
        Inches(4.15),
        Inches(3.9),
        "Supporting work from the draft",
        [
            "Complete system models and diagrams that guide implementation.",
            "Use Visual Paradigm and StarUML for structure, behavior, and interaction models.",
            "Prepare and clean selected datasets.",
            "Establish the evaluation setup for measuring performance.",
        ],
        SLATE_BG,
    )
    number(slide, slide_no)
    slide_no += 1

    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    title(slide, "Sitemap Showing User Interface Flow", "Existing visual artifact from the project folder", "UI")
    rect(slide, Inches(0.65), Inches(1.5), Inches(12.05), Inches(5.45), SLATE_BG, radius=True)
    image_box(slide, SITE_MAP, Inches(0.88), Inches(1.72), Inches(11.6), Inches(4.85))
    text(
        slide,
        Inches(0.9),
        Inches(6.63),
        Inches(11.0),
        Inches(0.25),
        "This flow captures the shared landing experience plus separate applicant and employer dashboards.",
        12,
        SLATE,
    )
    number(slide, slide_no)
    slide_no += 1

    slide = prs.slides.add_slide(blank)
    set_bg(slide, CREAM)
    title(slide, "Use Case Diagram", "User-facing actions already identified by the team", "Use Cases")
    rect(slide, Inches(0.65), Inches(1.5), Inches(8.1), Inches(5.45), SLATE_BG, radius=True)
    image_box(slide, USE_CASE, Inches(0.82), Inches(1.7), Inches(7.75), Inches(4.95))
    card(
        slide,
        Inches(9.0),
        Inches(1.85),
        Inches(3.45),
        Inches(3.9),
        "Key interactions shown",
        [
            "Upload a resume or CV.",
            "Upload a job description.",
            "View job lists and candidate lists.",
            "Check process status.",
            "View explanations and missing competencies.",
            "Inspect competency profiles as part of transparency.",
        ],
        PALE_GOLD,
    )
    text(
        slide,
        Inches(8.98),
        Inches(6.02),
        Inches(3.45),
        Inches(0.5),
        "These use cases support the functional requirements already drafted in the midsemester document.",
        11.5,
        NAVY,
    )
    number(slide, slide_no)

    prs.save(OUT_PPTX)
    return prs


def write_notes():
    notes = """# Job Hunt Midsemester Presentation Notes

## Slide 1 - Project title slide
- Project Name: Job Hunt
- Team Name: Free Labour
- Presented by:
  - Alwyn Sterling
  - Dominic Scott
  - Leah-Jay Holness
  - Ron-Hugh Walters
  - Tianna Bassaragh

## Slide 2 - Objectives Clearly Defined
### Background or context
We observe that modern recruitment is often inefficient and opaque. Employers must manually sift through many applications of varying relevance, while applicants often receive limited feedback about why they were not selected. This creates wasted effort on both sides and can cause qualified candidates to be missed. We are addressing this problem by formalising the comparison between job requirements and candidate competencies into a structured computational process.

We are not trying to model full hiring outcomes or replace recruiters. We are specifically scoping our project as an initial matching and recommendation stage, focused on whether the competencies expressed in a resume satisfy the requirements expressed in a job description. To keep the problem manageable and evaluation realistic, we are limiting our scope to computing related roles.

## Slide 3 - Solution
We aim to build a Job Matching and Qualification System for computing related roles. Our system will accept resumes from candidates and job descriptions from employers, extract relevant competencies from both, map them into a shared O*NET based competency space, and determine whether a candidate is qualified for a given role. Candidates who satisfy required criteria will then be ranked by suitability, while jobs may also be ranked for candidates by compatibility.

Our central objective is to reduce two major failure types: recommending candidates who are actually unqualified, and missing candidates who genuinely satisfy the requirements. In other words, we are trying to approximate the first pass judgement a skilled human recruiter makes in a more consistent and structured way.

## Slide 4 - Project scope
### What we will definitely do
- accept resumes and job descriptions as input
- extract competencies from unstructured text
- map extracted competencies to O*NET descriptors
- generate structured competency profiles for jobs and candidates
- determine whether required competencies are satisfied
- rank qualified candidates for employers and relevant jobs for candidates

### What we will definitely not do
We will not build a complete hiring platform. We will not cover later recruitment stages such as interviews, offer management, onboarding, or final hiring decisions. We will also not claim to predict actual job performance or guarantee real world hiring accuracy, since our system only operates on information available in resumes and job descriptions.

### What we might do
If time permits, we may include features that correspond to the non core requirements of the system. These include allowing employers to explicitly mark competencies as required or preferred, allowing users to adjust ranking weights within defined bounds, generating explanations for qualification decisions, producing competency gap reports for candidates who are not qualified, and allowing users to filter results by criteria to be decided later. These features are valuable, but they are secondary to ensuring that the system can first extract competencies reliably, determine qualification correctly, and produce useful rankings.

## Slide 5 - Technical issues and why the project is non-trivial
Our project is non trivial because the hardest part is not building a website, but converting messy, unstructured text into a reliable competency representation that can support qualification decisions. Resumes and job descriptions are often vague, inconsistent, and full of synonyms, implied skills, and uneven detail. Extracting competencies accurately and then mapping them into a fixed O*NET framework is technically demanding.

A second challenge is defining "qualified" in a defensible way. We are treating qualification as multi dimensional, based on factors such as skill coverage, relevance of experience and education, rather than a single raw score. This means we must combine multiple signals without making the system arbitrary or opaque.

A third challenge is that our approach assumes useful structure in a high dimensional competency space. We also have to consider the fact that linear relationships may not always hold, and that the system must scale beyond a tiny test dataset. This means the project involves real modelling, representation, and computational tradeoffs, not just basic CRUD functionality.

Finally, evaluation itself is difficult. Ground truth for "who is truly qualified" is rarely readily available, so we must define a credible testing strategy using O*NET occupational data, public datasets, and approximate labelled evaluation pairs. Our methodology therefore evaluates the system in two stages: first the binary qualification decision using precision, recall, and F1, and then ranking quality using metrics such as Precision@K, nDCG@K, and MRR if time permits.

## Slide 6 - Knowledge from the curriculum required
To solve this problem, we must draw on knowledge from several courses in the computing curriculum. COMP2171 and COMP2140 will support our documentation, requirements analysis, and system modelling, including the production of diagrams and structured design artefacts needed to communicate the system clearly. COMP3161 will support the database aspects of the project, including how we store resumes, job postings, extracted competencies, and generated results in a structured way.

We will also rely on COMP3162 Data Science for data cleaning, preprocessing, and potentially some of the modelling and evaluation work needed to handle messy resume and job description data. Since our implementation will make use of Python, COMP1127 provides foundational programming knowledge that will support the backend logic, text processing, and experimentation. INFO2180 will be relevant to the design and implementation of the front end through which employers and candidates interact with the system.

In addition, COMP3220 is relevant because the project involves intelligent decision making in the form of competency extraction, matching, and qualification assessment. COMP2211 will also be useful, particularly for thinking about optimization in ranking and decision processes, where we need to balance multiple factors in a structured way.

Taken together, these courses show that the project is not just a simple web application. It requires us to integrate knowledge from software engineering, databases, web development, artificial intelligence, data science, programming, and optimization, which is part of why we see it as a genuine capstone level problem.

## Slide 7 - Progress on target
### Work division
- Alwyn Sterling: Backend development and database implementation
- Dominic Scott: Backend development and database implementation
- Ron-Hugh Walters: Documentation, system modelling, data cleaning, and data parsing
- Leah-Jay Holness: User interface design
- Tianna Bassaragh: User interface design

## Slide 8 - What remains to be done
- Frontend: candidate and employer interfaces, page flow, and user interaction.
- Backend + Database: API endpoints, document processing flow, and storage for jobs, resumes, profiles, and matches.
- Implement the matching pipeline: resume parsing, job parsing, mapping, qualification, and ranking logic.
- Validate the system: prepare datasets, clean data, integrate evaluation, and test the pipeline.

## Slide 9 - Sitemap showing User Interface flow
Use the colored Real Site Map from the project folder.

## Slide 10 - Use Case Diagram
Use Case Diagram 7 from the project folder.
"""
    OUT_NOTES.write_text(notes, encoding="utf-8")


if __name__ == "__main__":
    ensure_dirs()
    build()
    write_notes()
    print(f"Created {OUT_PPTX}")
    print(f"Created {OUT_NOTES}")
