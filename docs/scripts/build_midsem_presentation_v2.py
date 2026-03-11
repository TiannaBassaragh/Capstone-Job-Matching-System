from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as ShapeType
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
DOCS = ROOT / "docs"
OUT_PPTX = DOCS / "job-hunt-midsemester-presentation.pptx"
OUT_NOTES = DOCS / "job-hunt-midsemester-notes.md"
SITE_MAP = DOCS / "Sitemap" / "Real Site Map.png"
USE_CASE = DOCS / "use_case_diagram" / "UseCaseDiagram7.png"

NAVY = RGBColor(17, 33, 62)
SLATE = RGBColor(45, 62, 84)
GOLD = RGBColor(224, 180, 74)
TEAL = RGBColor(54, 138, 127)
CREAM = RGBColor(247, 243, 235)
WHITE = RGBColor(255, 255, 255)
INK = RGBColor(31, 42, 55)
BLUE = RGBColor(225, 235, 247)
GREEN = RGBColor(228, 242, 236)
PALE_GOLD = RGBColor(247, 235, 199)
ROSE = RGBColor(246, 225, 220)
SLATE_BG = RGBColor(234, 238, 243)

TITLE_FONT = "Georgia"
BODY_FONT = "Aptos"


def style(paragraph, size, color, bold=False, font=BODY_FONT):
    if not paragraph.runs:
        paragraph.add_run()
    for run in paragraph.runs:
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, left, top, width, height, fill, radius=False):
    kind = ShapeType.ROUNDED_RECTANGLE if radius else ShapeType.RECTANGLE
    shape = slide.shapes.add_shape(kind, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    return shape


def text(slide, left, top, width, height, value, size=18, color=INK, bold=False, font=BODY_FONT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Pt(6)
    frame.margin_right = Pt(6)
    frame.margin_top = Pt(4)
    frame.margin_bottom = Pt(4)
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.text = value
    p.alignment = align
    style(p, size, color, bold=bold, font=font)
    return box


def bullets(slide, left, top, width, height, items, size=15, color=INK):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Pt(8)
    frame.margin_right = Pt(8)
    frame.margin_top = Pt(4)
    frame.margin_bottom = Pt(4)
    frame.clear()
    for i, item in enumerate(items):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = item
        p.bullet = True
        p.space_after = Pt(7)
        style(p, size, color)
    return box


def title(slide, value, subtitle, section):
    rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.34), GOLD)
    text(slide, Inches(0.7), Inches(0.55), Inches(9.7), Inches(0.72), value, 25, NAVY, True, TITLE_FONT)
    text(slide, Inches(0.72), Inches(1.18), Inches(9.7), Inches(0.36), subtitle, 11, TEAL, True)
    tag = rect(slide, Inches(10.62), Inches(0.62), Inches(1.95), Inches(0.36), TEAL, radius=True)
    p = tag.text_frame.paragraphs[0]
    p.text = section
    p.alignment = PP_ALIGN.CENTER
    style(p, 11, WHITE, True)


def card(slide, left, top, width, height, heading, items, fill):
    rect(slide, left, top, width, height, fill, radius=True)
    text(slide, left + Inches(0.12), top + Inches(0.1), width - Inches(0.24), Inches(0.34), heading, 14, NAVY, True, TITLE_FONT)
    bullets(slide, left + Inches(0.08), top + Inches(0.46), width - Inches(0.16), height - Inches(0.54), items, 11.5)


def step(slide, left, top, heading, body, fill):
    rect(slide, left, top, Inches(1.45), Inches(1.1), fill, radius=True)
    text(slide, left + Inches(0.1), top + Inches(0.08), Inches(1.25), Inches(0.28), heading, 13, NAVY, True, TITLE_FONT, PP_ALIGN.CENTER)
    text(slide, left + Inches(0.08), top + Inches(0.38), Inches(1.28), Inches(0.55), body, 10.5, INK, False, BODY_FONT, PP_ALIGN.CENTER)


def chevron(slide, left, top):
    shape = slide.shapes.add_shape(ShapeType.CHEVRON, left, top, Inches(0.42), Inches(0.54))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.color.rgb = GOLD


def image_box(slide, path, left, top, width, height):
    with Image.open(path) as image:
        iw, ih = image.size
    box_ratio = float(width) / float(height)
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        draw_width = width
        draw_height = int(width / img_ratio)
    else:
        draw_height = height
        draw_width = int(height * img_ratio)
    draw_left = left + int((width - draw_width) / 2)
    draw_top = top + int((height - draw_height) / 2)
    slide.shapes.add_picture(str(path), draw_left, draw_top, width=draw_width, height=draw_height)


def number(slide, value, dark=False):
    color = PALE_GOLD if dark else SLATE
    text(slide, Inches(12.52), Inches(7.0), Inches(0.4), Inches(0.2), str(value), 10, color, False, BODY_FONT, PP_ALIGN.RIGHT)


def add_title_slide(prs, blank, slide_no):
    slide = prs.slides.add_slide(blank)
    set_bg(slide, NAVY)
    rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.34), GOLD)
    rect(slide, Inches(0), Inches(6.92), Inches(13.333), Inches(0.58), TEAL)
    rect(slide, Inches(8.48), Inches(1.1), Inches(3.9), Inches(4.9), PALE_GOLD, radius=True)
    text(slide, Inches(0.78), Inches(0.98), Inches(4.3), Inches(0.3), "COMP3901 - Capstone Project", 14, PALE_GOLD, True)
    text(slide, Inches(0.78), Inches(1.45), Inches(5.8), Inches(0.9), "Job Hunt", 31, WHITE, True, TITLE_FONT)
    text(slide, Inches(0.8), Inches(2.28), Inches(6.2), Inches(0.9), "Midsemester Presentation", 20, WHITE, True)
    text(slide, Inches(0.8), Inches(3.18), Inches(6.35), Inches(1.05), "A web-based job matching and qualification system for computing-related roles.", 18, WHITE)
    text(slide, Inches(0.8), Inches(4.45), Inches(2.0), Inches(0.25), "Team Name", 12, PALE_GOLD, True)
    text(slide, Inches(0.8), Inches(4.7), Inches(2.5), Inches(0.45), "Free Labour", 22, WHITE, True, TITLE_FONT)
    text(slide, Inches(8.84), Inches(1.42), Inches(2.6), Inches(0.3), "Team Members", 12, NAVY, True)
    bullets(slide, Inches(8.82), Inches(1.8), Inches(3.0), Inches(3.8), ["Alwyn Sterling", "Dominic Scott", "Leah-Jay Holness", "Ron-Hugh Walters", "Tianna Bassaragh"], 14)
    number(slide, slide_no, dark=True)


def add_team_slide(prs, blank, slide_no):
    slide = prs.slides.add_slide(blank)
    set_bg(slide, CREAM)
    title(slide, "Team Introduction", "Free Labour and the project we are tackling", "Overview")
    text(slide, Inches(0.72), Inches(1.55), Inches(5.45), Inches(0.9), "We are building a web application that matches candidate resumes to computing jobs by extracting, mapping, qualifying, and ranking competencies.", 18, NAVY, True)
    card(slide, Inches(0.72), Inches(2.72), Inches(5.35), Inches(2.95), "Team and roles", ["Alwyn Sterling - backend development and database implementation", "Dominic Scott - backend development and database implementation", "Ron-Hugh Walters - documentation, system modelling, data cleaning, and data parsing", "Leah-Jay Holness - user interface design", "Tianna Bassaragh - user interface design"], BLUE)
    card(slide, Inches(6.45), Inches(1.72), Inches(5.9), Inches(1.45), "Project summary", ["Inputs: resumes from candidates and job descriptions from employers.", "Outputs: qualification decisions plus ranked results for both sides."], GREEN)
    card(slide, Inches(6.45), Inches(3.45), Inches(5.9), Inches(2.18), "Why this project matters", ["Recruitment is often inefficient and opaque.", "Qualified candidates can be overlooked.", "Employers spend time filtering weak matches.", "Applicants rarely get useful feedback."], PALE_GOLD)
    number(slide, slide_no)


def add_problem_slide(prs, blank, slide_no):
    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    title(slide, "Problem Background and Context", "Why this problem is worth solving", "Context")
    bullets(slide, Inches(0.72), Inches(1.55), Inches(6.3), Inches(4.9), ["Modern recruitment is inefficient for both employers and applicants.", "Job descriptions and resumes are messy, unstructured documents.", "Employers need a faster first-pass filter for candidate suitability.", "Applicants benefit from clearer, more transparent matching outcomes.", "The project focuses on computing-related roles to keep scope realistic."], 17)
    card(slide, Inches(7.35), Inches(1.8), Inches(5.1), Inches(1.4), "Problem framing", ["We are formalising comparison between job requirements and candidate competencies.", "The system acts as an initial matching and recommendation stage."], PALE_GOLD)
    card(slide, Inches(7.35), Inches(3.55), Inches(5.1), Inches(1.55), "Not the goal", ["This is not a full hiring platform.", "We are not predicting final hiring outcomes or job performance."], ROSE)
    number(slide, slide_no)


def add_solution_slide(prs, blank, slide_no):
    slide = prs.slides.add_slide(blank)
    set_bg(slide, CREAM)
    title(slide, "Solution Overview", "How the web application is intended to work", "Solution")
    bullets(slide, Inches(0.72), Inches(1.55), Inches(5.5), Inches(4.9), ["Candidates upload resumes and employers submit job descriptions.", "The system extracts competencies from both sources.", "Extracted content is mapped into an O*NET-based competency space.", "The system checks whether required competencies are satisfied.", "Qualified matches are ranked for employers and for candidates."], 17)
    step(slide, Inches(6.45), Inches(1.9), "Input", "Resume and job description submission", PALE_GOLD)
    chevron(slide, Inches(7.93), Inches(2.18))
    step(slide, Inches(8.42), Inches(1.9), "Extract", "Competency extraction from text", GREEN)
    chevron(slide, Inches(9.9), Inches(2.18))
    step(slide, Inches(10.39), Inches(1.9), "Map", "O*NET-based representation", BLUE)
    chevron(slide, Inches(8.7), Inches(3.95))
    step(slide, Inches(8.2), Inches(4.28), "Qualify", "Check required competencies", ROSE)
    chevron(slide, Inches(9.68), Inches(4.56))
    step(slide, Inches(10.17), Inches(4.28), "Rank", "Order the strongest matches", PALE_GOLD)
    number(slide, slide_no)


def add_scope_slide(prs, blank, slide_no):
    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    title(slide, "Objectives Clearly Defined", "Core deliverables, boundaries, and stretch goals", "Scope")
    card(slide, Inches(0.62), Inches(1.55), Inches(4.05), Inches(4.95), "Core deliverables", ["Accept resumes and job descriptions.", "Extract competencies from unstructured text.", "Map extracted competencies to O*NET descriptors.", "Generate structured competency profiles.", "Determine whether required competencies are satisfied.", "Rank qualified candidates and relevant jobs."], GREEN)
    card(slide, Inches(4.9), Inches(1.55), Inches(4.0), Inches(4.95), "Clearly out of scope", ["A full hiring platform.", "Interviews, offers, onboarding, or final hiring decisions.", "Predicting job performance.", "Claiming guaranteed real-world hiring accuracy."], ROSE)
    card(slide, Inches(9.1), Inches(1.55), Inches(3.62), Inches(4.95), "Stretch goals", ["Mark competencies as required or preferred.", "Adjust ranking weights within defined bounds.", "Generate qualification explanations.", "Produce competency gap reports.", "Filter results by later decided criteria."], PALE_GOLD)
    number(slide, slide_no)


def add_requirements_slide(prs, blank, slide_no):
    slide = prs.slides.add_slide(blank)
    set_bg(slide, CREAM)
    title(slide, "Functional Requirements Overview", "The most important behaviours of the system", "Requirements")
    card(slide, Inches(0.72), Inches(1.65), Inches(3.0), Inches(1.8), "Input and storage", ["Upload or enter resumes and job descriptions.", "Store structured profiles for later matching."], BLUE)
    card(slide, Inches(3.95), Inches(1.65), Inches(3.0), Inches(1.8), "Competency profiling", ["Extract competencies from unstructured text.", "Map extracted signals to O*NET descriptors."], GREEN)
    card(slide, Inches(7.18), Inches(1.65), Inches(2.95), Inches(1.8), "Qualification logic", ["Determine whether required competencies are satisfied.", "Separate qualified from not qualified cases."], PALE_GOLD)
    card(slide, Inches(10.36), Inches(1.65), Inches(2.2), Inches(1.8), "Ranked results", ["Rank candidates for employers.", "Rank jobs for candidates."], ROSE)
    card(slide, Inches(1.45), Inches(4.08), Inches(5.05), Inches(1.7), "Important supporting behaviour", ["Show the most useful results, not just raw matches.", "Keep the design realistic for a web application workflow."], SLATE_BG)
    card(slide, Inches(6.82), Inches(4.08), Inches(5.1), Inches(1.7), "Stretch features stay separate", ["Explanation, weighting controls, preferred vs required tags, and gap reports are optional additions.", "They are not presented as guaranteed deliverables."], BLUE)
    number(slide, slide_no)


def add_method_slide(prs, blank, slide_no):
    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    title(slide, "System Flow and Methodology", "From document input to ranked recommendations", "Method")
    step(slide, Inches(0.9), Inches(2.0), "1. Collect", "Candidate resumes and employer job descriptions", PALE_GOLD)
    chevron(slide, Inches(2.38), Inches(2.28))
    step(slide, Inches(2.88), Inches(2.0), "2. Extract", "Find competency signals in unstructured text", GREEN)
    chevron(slide, Inches(4.36), Inches(2.28))
    step(slide, Inches(4.86), Inches(2.0), "3. Map", "Represent extracted signals with O*NET descriptors", BLUE)
    chevron(slide, Inches(6.34), Inches(2.28))
    step(slide, Inches(6.84), Inches(2.0), "4. Qualify", "Check whether required competencies are satisfied", ROSE)
    chevron(slide, Inches(8.32), Inches(2.28))
    step(slide, Inches(8.82), Inches(2.0), "5. Rank", "Order the strongest matches for both user groups", PALE_GOLD)
    rect(slide, Inches(1.15), Inches(4.45), Inches(10.9), Inches(1.25), NAVY, radius=True)
    text(slide, Inches(1.42), Inches(4.68), Inches(3.0), Inches(0.25), "Evaluation focus", 12, PALE_GOLD, True)
    bullets(slide, Inches(1.4), Inches(4.95), Inches(10.25), Inches(0.55), ["Qualification quality: precision, recall, and F1.", "Ranking quality if time permits: Precision@K, nDCG@K, and MRR."], 14, WHITE)
    number(slide, slide_no)


def add_sitemap_slide(prs, blank, slide_no):
    slide = prs.slides.add_slide(blank)
    set_bg(slide, CREAM)
    title(slide, "Sitemap", "Public, candidate, and employer areas of the web application", "Diagram")
    rect(slide, Inches(0.62), Inches(1.55), Inches(8.65), Inches(5.25), SLATE_BG, radius=True)
    image_box(slide, SITE_MAP, Inches(0.85), Inches(1.78), Inches(8.2), Inches(4.75))
    card(slide, Inches(9.55), Inches(1.75), Inches(2.95), Inches(3.3), "How to read it", ["Public area: landing page, about, contact, sign up, and log in.", "Candidate area: resume upload, matches, notifications, settings, and match detail.", "Employer area: job posting, candidate matches, notifications, settings, and match detail."], PALE_GOLD)
    text(slide, Inches(0.88), Inches(6.56), Inches(8.1), Inches(0.25), "Caption: The sitemap shows the main public pages plus the candidate and employer dashboard paths.", 11.5, SLATE)
    number(slide, slide_no)


def add_use_case_slide(prs, blank, slide_no):
    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    title(slide, "Use Case Diagram", "Main actors and interactions supported by the system", "Diagram")
    rect(slide, Inches(0.62), Inches(1.55), Inches(8.65), Inches(5.25), SLATE_BG, radius=True)
    image_box(slide, USE_CASE, Inches(0.85), Inches(1.78), Inches(8.2), Inches(4.75))
    card(slide, Inches(9.55), Inches(1.75), Inches(2.95), Inches(3.4), "Actors and interactions", ["Actors interact with resume upload and job description upload flows.", "Users view job lists, candidate lists, and profile editing actions.", "The model also includes viewing details and explanations for transparency."], BLUE)
    text(slide, Inches(0.88), Inches(6.56), Inches(8.1), Inches(0.25), "Caption: The use case diagram highlights the core user actions around submission, viewing, and explanation.", 11.5, SLATE)
    number(slide, slide_no)


def add_challenges_slide(prs, blank, slide_no):
    slide = prs.slides.add_slide(blank)
    set_bg(slide, CREAM)
    title(slide, "Technical Challenges and Why the Project Is Non-Trivial", "This is a modelling problem, not just a CRUD app", "Challenge")
    card(slide, Inches(0.72), Inches(1.7), Inches(2.85), Inches(2.0), "Extraction", ["Competencies must be pulled from messy resumes and job descriptions.", "The source documents are inconsistent and incomplete."], BLUE)
    card(slide, Inches(3.78), Inches(1.7), Inches(2.85), Inches(2.0), "O*NET mapping", ["Extracted text must be mapped into a shared competency space.", "That mapping has to be meaningful enough for comparison."], GREEN)
    card(slide, Inches(6.84), Inches(1.7), Inches(2.85), Inches(2.0), "Defensible qualification", ["Qualification must reflect skill coverage, experience, and education.", "The decision cannot collapse into an arbitrary single score."], PALE_GOLD)
    card(slide, Inches(9.9), Inches(1.7), Inches(2.75), Inches(2.0), "Evaluation", ["Ground truth is limited, so evaluation must be careful and realistic.", "Metrics include precision, recall, F1, and ranking metrics if time permits."], ROSE)
    rect(slide, Inches(0.72), Inches(4.28), Inches(11.9), Inches(1.35), NAVY, radius=True)
    text(slide, Inches(0.98), Inches(4.52), Inches(5.0), Inches(0.25), "Why this is capstone-level work", 12, PALE_GOLD, True)
    bullets(slide, Inches(0.95), Inches(4.8), Inches(11.3), Inches(0.5), ["It combines web development, data representation, AI-style matching, evaluation design, and scalability concerns."], 14, WHITE)
    number(slide, slide_no)


def add_progress_slide(prs, blank, slide_no):
    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    title(slide, "Progress on Target and Well Organised", "What has been done, what remains, and how work is divided", "Progress")
    card(slide, Inches(0.62), Inches(1.68), Inches(3.95), Inches(4.95), "What has been done", ["Project problem and scope have been defined.", "Core methodology and O*NET-based approach have been outlined.", "Functional requirements have been drafted.", "Sitemap and use case diagrams have been produced.", "Team roles and work areas are clearly assigned."], BLUE)
    card(slide, Inches(4.7), Inches(1.68), Inches(3.95), Inches(4.95), "What remains", ["Build the frontend workflow for candidates and employers.", "Complete backend and database integration.", "Implement the matching pipeline end to end.", "Prepare datasets and run validation.", "Connect evaluation results to the system story."], PALE_GOLD)
    card(slide, Inches(8.78), Inches(1.68), Inches(3.95), Inches(4.95), "Work division", ["Alwyn Sterling - backend and database", "Dominic Scott - backend and database", "Ron-Hugh Walters - documentation, modelling, data work", "Leah-Jay Holness - UI design", "Tianna Bassaragh - UI design"], GREEN)
    number(slide, slide_no)


def add_curriculum_slide(prs, blank, slide_no):
    slide = prs.slides.add_slide(blank)
    set_bg(slide, CREAM)
    title(slide, "Knowledge From the Curriculum and Next Steps", "How the project draws on coursework and where it goes next", "Close")
    card(slide, Inches(0.72), Inches(1.65), Inches(6.05), Inches(4.9), "Curriculum knowledge used", ["COMP2171 and COMP2140 - documentation and system modelling", "COMP3161 - database design and storage", "COMP3162 - data cleaning and possible modelling", "INFO2180 - web development and frontend design", "COMP3220 - intelligent matching concepts", "COMP1127 - Python implementation knowledge", "COMP2211 - optimization ideas"], BLUE)
    card(slide, Inches(7.02), Inches(1.65), Inches(5.15), Inches(2.05), "Immediate next steps", ["Implement the frontend and backend workflow.", "Build the O*NET-based matching pipeline.", "Prepare datasets for testing and evaluation."], PALE_GOLD)
    card(slide, Inches(7.02), Inches(4.02), Inches(5.15), Inches(1.85), "Midsemester takeaway", ["The project is clearly scoped, technically non-trivial, and organized around a realistic implementation plan."], ROSE)
    number(slide, slide_no)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slide_no = 1

    add_title_slide(prs, blank, slide_no); slide_no += 1
    add_team_slide(prs, blank, slide_no); slide_no += 1
    add_problem_slide(prs, blank, slide_no); slide_no += 1
    add_solution_slide(prs, blank, slide_no); slide_no += 1
    add_scope_slide(prs, blank, slide_no); slide_no += 1
    add_requirements_slide(prs, blank, slide_no); slide_no += 1
    add_method_slide(prs, blank, slide_no); slide_no += 1
    add_sitemap_slide(prs, blank, slide_no); slide_no += 1
    add_use_case_slide(prs, blank, slide_no); slide_no += 1
    add_challenges_slide(prs, blank, slide_no); slide_no += 1
    add_progress_slide(prs, blank, slide_no); slide_no += 1
    add_curriculum_slide(prs, blank, slide_no)

    prs.save(OUT_PPTX)


def write_notes():
    notes = """# Job Hunt Midsemester Presentation Notes

## Slide 1 - Title slide
- Project name: Job Hunt
- Team name: Free Labour
- Course: COMP3901 - Capstone Project
- Midsemester presentation for a web-based job matching and qualification system.

## Slide 2 - Team introduction
- One-sentence summary: we are building a web application that extracts, maps, qualifies, and ranks competencies from resumes and job descriptions.
- Team roles:
  - Alwyn Sterling: backend development and database implementation
  - Dominic Scott: backend development and database implementation
  - Ron-Hugh Walters: documentation, system modelling, data cleaning, and data parsing
  - Leah-Jay Holness: user interface design
  - Tianna Bassaragh: user interface design

## Slide 3 - Problem background and context
- Recruitment is inefficient and often opaque.
- Employers need better first-pass filtering.
- Applicants need clearer and more transparent matching outcomes.
- Scope is limited to computing-related roles.
- This is not a full hiring platform.

## Slide 4 - Solution overview
- Inputs are resumes from candidates and job descriptions from employers.
- Competencies are extracted from both sources.
- Extracted content is mapped into an O*NET-based competency space.
- The system determines whether required competencies are satisfied.
- Qualified matches are ranked for employers and for candidates.

## Slide 5 - Objectives clearly defined
- Core deliverables are the must-do features.
- Out-of-scope items include later hiring stages and job-performance prediction.
- Stretch goals are should-level features only and should not be overpromised.

## Slide 6 - Functional requirements overview
- Focus on the most important behaviours: input, competency profiling, qualification, and ranked results.
- Supportive features are framed as optional rather than guaranteed.

## Slide 7 - System flow and methodology
- Pipeline: collect, extract, map, qualify, and rank.
- Evaluation focus: precision, recall, F1, and ranking metrics if time permits.

## Slide 8 - Sitemap
- Explain the public pages first.
- Then explain the candidate dashboard flow.
- Then explain the employer dashboard flow.

## Slide 9 - Use case diagram
- Point out the main actors and the core interactions.
- Emphasize submission, viewing, and explanation-related actions.

## Slide 10 - Technical challenges
- Extraction from messy text.
- Mapping into O*NET descriptors.
- Defining qualification in a defensible multi-dimensional way.
- Evaluating quality with limited ground truth.

## Slide 11 - Progress on target and well organised
- Done: scope, methodology, requirements, and diagrams.
- Remaining: frontend, backend/database, pipeline, datasets, and validation.
- The work division is already clearly assigned.

## Slide 12 - Knowledge from the curriculum and next steps
- Highlight how multiple courses feed into the project.
- Close on the next implementation steps and the realism of the current plan.
"""
    OUT_NOTES.write_text(notes, encoding="utf-8")


if __name__ == "__main__":
    build()
    write_notes()
    print(f"Created {OUT_PPTX}")
    print(f"Created {OUT_NOTES}")
